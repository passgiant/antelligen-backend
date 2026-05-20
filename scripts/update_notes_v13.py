"""v1.3 PPTX 발표자 노트 업데이트 스크립트."""
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "docs" / "angelligen" / "antelligen_presentation_ver1.3.pptx"

PART_NOTE = "잠시 다음 파트로 넘어가겠습니다."
CONT_NOTE = ""

NOTES: dict[int, str] = {
    1: (
        "안녕하세요, Antelligen 팀입니다. 7인 팀으로 AI 멀티 에이전트 기반 투자 인텔리전스 플랫폼을 개발했습니다.\n"
        "핵심 메시지는 AI가 대신 투자하는 시대가 아니라, AI의 판단을 사람이 검증하고 통제하는 구조입니다.\n"
        "발표 순서는 문제 인식 → 해결 방안 → 성장 전략 → 아키텍처 → 기능 소개 → 팀원별 개발 내용 → 시연 순으로 진행하겠습니다."
    ),
    2: PART_NOTE,
    3: (
        "왜 멀티 에이전트인가, 가장 근본적인 질문입니다.\n"
        "단일 LLM은 데이터 수집·분석·의사결정을 하나의 모델에 몰아넣어야 하는 구조적 한계가 있습니다.\n"
        "현실의 투자 판단은 여러 전문가의 협업으로 이루어지는데, 기존 AI 시스템은 이를 반영하지 못했습니다.\n"
        "Antelligen 5가지 목표: 구조적 AI 시스템 구축 / 확장 가능한 아키텍처 / 현실적 협업 흐름 반영 / 사용자 맞춤 분석 / 자동화·효율성 향상."
    ),
    4: (
        "TAM은 글로벌 AI 에이전트 시장 — 2025년 76억 달러에서 2033년 1,820억 달러로 연평균 50% 성장이 예상됩니다.\n"
        "SAM은 국내 AI 시장 54.7억 달러(2024) → 538.7억 달러(2032), 그중 BFSI(금융·보험)의 AI 도입 비중이 가장 높습니다.\n"
        "SOM은 중소기업과 개인 투자자를 대상으로, 클라우드 기반 AI 웹서비스 형태로 공략합니다.\n"
        "블룸버그 월 200만 원+에 비해 합리적 가격으로 한국 공시 중심 서비스를 제공하는 것이 저희 포지셔닝입니다."
    ),
    5: (
        "멀티 에이전트 시스템 도입 시 고객이 겪는 3대 문제입니다.\n"
        "첫째, 가시성 문제 — 에이전트가 많아질수록 어느 단계에서 실패했는지 추적이 어렵습니다. LangSmith Tracing으로 해결했습니다.\n"
        "둘째, 비용 통제 — 토큰 비용이 통제 없이 급증하는 Runaway Cost 문제. 다층 캐시와 조건부 Startup으로 대응합니다.\n"
        "셋째, 거버넌스 — 금융 데이터의 민감성, Prompt Injection, Audit Trail 요구. 설계 단계부터 내재화했습니다.\n"
        "핵심 인사이트: 신뢰와 통제가 없는 AI는 금융에서 실패한다."
    ),
    6: PART_NOTE,
    7: (
        "데이터 신뢰도·커버리지 문제 5가지를 해결했습니다.\n"
        "P-1: 뉴스 분석이 하드코딩 8종목 한정 → TickerKeywordResolverPort로 전 KRX 2,500+ 종목 지원.\n"
        "P-2: 미국 주식 미지원 → MarketRegion 추상화로 SerpAPI(뉴스)·yfinance(재무)·SEC EDGAR(공시) 경로 분기.\n"
        "P-3: 잠정실적 누락 → OpenDartPreliminaryEarningsProvider로 어닝 갭 윈도우 포착.\n"
        "P-4: 출처 무관 단순 평균 → SourceTier 4단계(HIGH/MEDIUM/MEDIUM_LOW/LOW) 가중 평균 적용.\n"
        "P-5: 엔터주 SNS 시그널 묻힘 → 섹터별 소스 가중치 오버라이드로 팬덤 반응을 시그널에 반영."
    ),
    8: (
        "운영·비용·UX 문제 6가지 해결입니다.\n"
        "P-6: LLM·YouTube quota 폭주 → 4h/7d/1d/1h 다층 캐시 + 조건부 Startup 잡으로 불필요한 호출 차단.\n"
        "P-7: 동시 요청 → DART 중복 호출 → Redis 분산 락(SET NX EX)으로 atomic 단일 호출 보장.\n"
        "P-8: 응답 지연 5~10초 → 1h Postgres 캐시 + 3-에이전트 병렬 + 사업개요 별도 async task.\n"
        "P-9: 시그널만 있고 맥락 없음 → BusinessOverview 카드로 RAG 기반 LLM 사업 요약 제공.\n"
        "P-10: 1개 에이전트 실패 → 전체 실패 → asyncio.gather(return_exceptions=True)로 부분 실패 허용.\n"
        "P-11: 왜 이 결론? → LangSmith Tracing + 출처 영속화 + Explainable 가중치로 투명성 확보.\n"
        "6대 해결 원칙 중 부분 실패는 기능과 거버넌스는 사전 설계가 설계 철학을 잘 요약합니다."
    ),
    9: (
        "기술 프레임워크 경쟁자로는 CrewAI, AutoGen, LangGraph, LlamaIndex가 있습니다.\n"
        "서비스형 플랫폼으로는 Lindy, AgentGPT, 기존 금융 알고리즘 도구들이 있습니다.\n"
        "시장에 AI 빌더와 트레이딩 도구는 많지만, 검증 가능한 의사결정 과정을 제공하는 서비스는 아직 드뭅니다.\n"
        "저희가 채우려는 빈자리가 바로 거기입니다."
    ),
    10: (
        "차별화 3축입니다.\n"
        "첫째, 신뢰성·거버넌스 — LangSmith Tracing으로 모든 LLM 입출력과 latency 추적, Source Tier로 출처 품질 수치화, Audit Trail로 분석 결과를 DB에 영속화.\n"
        "둘째, 고도화된 데이터 통합 — DART 사업보고서 RAG, 멀티 에이전트 5종, 한국 공시+미국 SEC+yfinance+잠정실적 통합.\n"
        "셋째, UX 차별화 — 사업개요 카드로 이 회사의 본질을 즉시 파악, 부분 실패도 UI에 표시해 신뢰성 유지.\n"
        "한마디로: 투자 결과가 아니라, 검증 가능한 의사결정 과정을 제공하는 서비스입니다."
    ),
    11: PART_NOTE,
    12: (
        "1차 타겟은 이 종목이 오늘 왜 떨어졌지?를 검색하는 개인 투자자입니다. 차트+사건+원인을 한 화면에서 보여줍니다.\n"
        "2차 타겟은 신입 애널리스트, 투자 콘텐츠 크리에이터, 학생 투자 동아리입니다.\n"
        "진입 채널 3가지: SEO(종목별 왜 떨어졌나 자동 페이지), 커뮤니티(디스코드·네이버 종목 카페), 인플루언서(투자 유튜버 협업).\n"
        "서비스 확장 경로는 웹사이트 → PWA → 모바일 앱 순서로 계획하고 있습니다."
    ),
    13: (
        "KPI를 4축으로 관리합니다.\n"
        "사용량: 초기 MAU 30~50명 → 확장 10,000명 / 7일 재방문율 30% → 40%.\n"
        "품질: 사용자 만족도 4.0 → 4.5 / 자신 있는 분석 비율 30% → 45%.\n"
        "비용 효율: 응답 속도 8초 → 3초 / 종목당 분석 비용 200원 → 50원. 현재 캐시와 분산 락으로 이 방향으로 나아가고 있습니다.\n"
        "사업성: 유료 전환율 3%, 월 매출 100만 원이 1차 마일스톤입니다."
    ),
    14: (
        "확장은 4방향입니다.\n"
        "시장·자산 확장: 일본·유럽·신흥국 / 암호화폐·원자재·채권·리츠 / 대안 데이터(위성·카드 소비).\n"
        "기능 확장: 포트폴리오 분석 / 실시간 알림 / 백테스트 / 일일·주간 리포트.\n"
        "수익 모델: 무료 → 개인 Pro 9,900원 → 프리미엄 29,000원 → API 과금 → 기업 전용 최대 5,000만 원.\n"
        "사용자 행동 → 검증 사례 → 규칙 개선 → 정확도 상승 → 재방문 플라이휠이 형성되면 1년 후 후발주자가 따라잡기 어려운 구조가 됩니다."
    ),
    15: PART_NOTE,
    16: (
        "기술 스택 한눈에 보겠습니다.\n"
        "백엔드: Python 3.13 / FastAPI / PostgreSQL 16(asyncpg+pgvector) / SQLAlchemy 2.0 async / Redis / APScheduler(20+ 정기 작업) / Docker Compose.\n"
        "AI/LLM: OpenAI(GPT-4o, gpt-5-mini, o1) / LangChain·LangGraph·LangSmith.\n"
        "데이터: yfinance·pykrx·FRED·Finnhub·GDELT / DART·SEC EDGAR·SerpAPI·YouTube Data API / kiwipiepy·BeautifulSoup4.\n"
        "프론트엔드: Next.js 16(App Router, Turbopack) / React 19 / TypeScript 5.9 / Jotai 2 / Tailwind CSS v4 / lightweight-charts·recharts."
    ),
    17: (
        "헥사고날 + DDD 아키텍처입니다.\n"
        "Clean Architecture 핵심 원칙 — 의존성은 항상 바깥에서 안쪽(도메인)으로만 흐릅니다.\n"
        "Domain 레이어는 순수 Python만 허용합니다. FastAPI·SQLAlchemy·Redis·Pydantic import 절대 금지.\n"
        "Application(UseCase) 레이어는 Port 인터페이스를 통해서만 외부 시스템에 접근합니다.\n"
        "Inbound Adapter(Router)는 비즈니스 로직 없이 Request → UseCase → Response 흐름만 담당합니다.\n"
        "DB나 외부 API가 바뀌어도 도메인 로직은 전혀 수정할 필요가 없습니다."
    ),
    18: (
        "24개 도메인으로 구성됩니다.\n"
        "사용자·인증 4개 / 시장 데이터 5개 / 콘텐츠·분석 7개 / 통합 에이전트(메인 오케스트레이터) 1개.\n"
        "요청 처리 흐름: HTTP → FastAPI Router → UseCase → Port → Repository(PostgreSQL)/External Client/Cache(Redis) → Domain Entity → Response DTO → HTTP Response.\n"
        "백그라운드 스케줄러 주요 4가지: 뉴스 수집(매일 06:00) / KRX 순매수(매 영업일 16:30) / SEC 13F(분기) / DART 포트폴리오(매월 1일)."
    ),
    19: CONT_NOTE,
    20: PART_NOTE,
    21: (
        "멀티 에이전트 오케스트레이션 구조입니다.\n"
        "메인 에이전트(오케스트레이터)가 asyncio.gather로 5개 서브 에이전트를 동시에 병렬 호출합니다.\n"
        "News Agent(SerpAPI+LLM) / Disclosure Agent(DART+RAG+LLM) / Finance Agent(yfinance+DART) / Sentiment Agent(SNS+GPT) / Macro Agent(YouTube+LLM).\n"
        "SourceTier 가중치 적용 → LangGraph로 통합 시그널 합성 → LangSmith Tracing → Bullish/Neutral/Bearish + Confidence 반환.\n"
        "부분 실패 허용(return_exceptions=True), Postgres 1h + Redis 7d 이중 캐시로 성능과 비용을 균형 잡았습니다."
    ),
    22: (
        "Source Tier 가중치 시스템입니다.\n"
        "HIGH(×1.0): DART·SEC·1군 IB 공식 리포트·기업 IR — 가장 신뢰도 높은 공식 출처.\n"
        "MEDIUM(×0.7): Bloomberg·Reuters·WSJ·FT·한경·매경·월가 애널리스트.\n"
        "MEDIUM_LOW(×0.5): 국내 IB 공식 리포트 — buy-bias 보정 적용.\n"
        "LOW(×0.3): SNS·일반 뉴스·커뮤니티.\n"
        "섹터 오버라이드 — HYBE·SM·JYP·YG 같은 엔터주 분석 시 YouTube·X·Instagram이 LOW에서 MEDIUM으로 자동 승격됩니다. 팬덤 반응이 실제 주가에 직접 영향을 미치는 특성을 반영한 것입니다."
    ),
    23: (
        "캐시·동시성·복원력 세 가지입니다.\n"
        "다층 캐시 — TTL은 도메인 진실성의 코드 표현입니다. 거시 스냅샷 4h / 회사 프로필 1d / 사업 개요(LLM) 7d / 통합 분석 1h / 회사 마스터 쿨다운 24h.\n"
        "Redis 분산 락 — SET key 1 NX EX 60으로 atomic single-shot. race condition 제로, TTL 60초로 데드락 자동 방지, Redis 장애 시에도 graceful degrade.\n"
        "부분 실패 허용 — asyncio.gather(return_exceptions=True)로 1개 에이전트가 죽어도 나머지로 시그널을 산출합니다. 부분 실패는 기능이라는 원칙의 코드 구현입니다."
    ),
    24: PART_NOTE,
    25: (
        "주요 기능 첫 번째 파트입니다.\n"
        "종합 종목 분석(/agent/query): 한 번의 요청으로 5개 에이전트 병렬 분석 → Bullish/Neutral/Bearish+Confidence, 사업개요 카드, Source Tier 분포 시각화. 한국·미국 동일 UI/UX.\n"
        "히스토리·인과관계 대시보드(/dashboard): 일봉 차트 + 5중 이상치 탐지 + LLM 원인 추론(왜 이 날 떨어졌는가) + 가설 신뢰도 + 거시 이벤트 매칭 + 마커 6종.\n"
        "Sentiment 분석(/sentiment): Reddit+네이버 종목토론 SNS 감정 점수화, GPT-5-mini 감정 분류, 밈 티커 부스트(×1.35), 부정 비율 vs VIX 비교.\n"
        "거시경제 현황판(/macro): 금리/유가/환율/경제 일정 자동 수집, Risk-ON/OFF 자동 판단+LLM 근거, 4h Redis 캐시로 실시간성과 비용 균형."
    ),
    26: (
        "주요 기능 두 번째 파트입니다.\n"
        "Smart Money(/smart-money): 외국인/기관/개인 순매수 랭킹, 집중 매수 종목 카드(3/5/10일 기간 선택), 종목 클릭 시 최근 30일 투자자별 추이 차트.\n"
        "글로벌 포트폴리오(/smart-money/global-portfolio): Buffett·Burry·Cathie Wood 등 19명 SEC 13F + 국민연금·미래에셋·박현주 등 10명 DART, 신규편입/비중확대/비중축소/청산 배지.\n"
        "관심종목+개인화(/settings/watchlist): 테마별 토글 등록 → 뉴스 자동 필터링 → 유튜브 병렬 검색 자동화.\n"
        "회사 프로필(/company-profile/{ticker}): DART+SEC 회사 마스터, 사업보고서 RAG→LLM 요약(7일 캐시), ETF/INDEX도 동일 카드 형식.\n"
        "카카오 OAuth: temp_token/user_token 분리로 가입 전 사용자도 핵심 기능 시연 가능."
    ),
    27: PART_NOTE,
    28: (
        "Antelligen 5인 개발팀입니다.\n"
        "김영진: 통합 에이전트·데이터 인프라·신뢰도 시스템 (Python·DDD·Hexagonal·LangChain·Redis).\n"
        "김세진: Sentiment 에이전트 전체·미국 주식 뉴스 보강 (DDD·BeautifulSoup4·OpenAI·Next.js·Jotai).\n"
        "문우성: 거시경제 현황판·시스템 아키텍처 정립 (Hexagonal·LLM·YouTube API).\n"
        "김민규: History Agent·Causality Agent·Dashboard (LangGraph·Next.js·lightweight-charts).\n"
        "이민영: 개인화 시스템·Smart Money 대시보드 (Next.js·React·Playwright·Jotai·DDD).\n"
        "이지은·정우진은 외부 피드백으로 팀을 지원해주셨습니다."
    ),
    29: CONT_NOTE,
    30: (
        "김영진 첫 번째 슬라이드 — 전 KRX 종목 지원과 미국 주식 확대입니다.\n"
        "Before: 뉴스 분석 가능 종목 하드코딩 8개, 미국 종목 미지원, 재무는 DART 정기공시만.\n"
        "After: TickerKeywordResolverPort로 DB에서 stock_name + alias 동적 조회. 000270(기아) → [기아, 기아차, KIA] 자동 생성. 전 KRX 2,500+ 종목 정상 동작.\n"
        "MarketRegion Enum으로 ticker 형식만 보고 KR/US 자동 추론합니다. 6자리 숫자 → KR, 1~5 알파벳 → US_NASDAQ.\n"
        "미국 경로: 뉴스는 SerpAPI(gl=us, hl=en), 재무는 yfinance 분기 EPS, 공시는 SEC EDGAR 8-K/10-K/10-Q."
    ),
    31: CONT_NOTE,
    32: (
        "김영진 두 번째 슬라이드 — Source Tier 4단계 가중치 시스템입니다.\n"
        "SourceTier를 Enum으로 정의하고 _DEFAULT_WEIGHTS로 HIGH:1.0 / MEDIUM:0.7 / MEDIUM_LOW:0.5 / LOW:0.3 매핑.\n"
        "엔터테인먼트 섹터 오버라이드: HYBE·SM·JYP·YG 분석 시 YouTube·X·Instagram → LOW → MEDIUM 자동 승격.\n"
        "종합 시그널 산출: 에이전트별 confidence × tier_weight 곱한 값을 합산, confidence_total로 나눈 avg_score가 0.2 초과면 bullish, -0.2 미만이면 bearish.\n"
        "가중치 수치는 환경 변수로 노출되어 있어 운영 중에도 실험적 튜닝이 가능합니다."
    ),
    33: CONT_NOTE,
    34: (
        "김영진 세 번째 슬라이드 — On-demand 분산 락, 조건부 Startup 잡, 사업개요 카드입니다.\n"
        "On-demand 분산 락: 동일 종목 동시 5개 요청이 들어와도 DART API는 1회만 호출. 첫 번째 요청이 Redis SET NX EX로 락 획득 → DART fetch → DB upsert → 락 해제. 나머지 요청은 락 대기 후 DB 결과 공유.\n"
        "조건부 Startup 잡: refresh_company_list는 24h 이내 성공 이력 있으면 스킵, refresh_market_risk는 Redis 4h 스냅샷 복원 → YouTube quota + LLM 비용 최대 절약.\n"
        "사업개요 카드: asyncio.create_task로 병렬 fetch(latency 거의 0), DART 사업보고서 RAG(최대 5청크·3,000자) → LLM 요약 → Redis 7일 캐시. LLM 컨텍스트 주입과 UI 카드 DTO 두 가지로 활용합니다."
    ),
    35: (
        "김세진 첫 번째 슬라이드 — Sentiment 도메인 신설입니다.\n"
        "담당 범위: 백엔드 21개 파일 신규(+1,727줄) / 프론트엔드 5개 파일 신규(+592줄).\n"
        "DDD 헥사고날 구조를 처음부터 적용해 adapter/application/domain 3계층으로 설계했습니다.\n"
        "SNS 수집기 3종: Reddit(.json 무인증 엔드포인트), 네이버 종목토론(BeautifulSoup4 HTML 스크래핑), 토스 커뮤니티(stub — SPA 봇 탐지로 Playwright 차기 이터레이션).\n"
        "Redis TTL 10분 캐시를 라우터 레벨에 배치해 UseCase 순수성(DDD 원칙)을 유지했습니다."
    ),
    36: CONT_NOTE,
    37: (
        "김세진 두 번째 슬라이드 — GPT 감정분석 + 밈/섹터 가중치 + DB 설계입니다.\n"
        "GPT 감정분석: 게시물 제목·본문 → GPT-5-mini(팀 표준 모델) → positive/negative/neutral 분류 → 종합 신호(bullish/bearish/neutral) + confidence + overall_negative_ratio(VIX 비교용).\n"
        "밈 티커 부스트: TSLA·GME·AMC·NVDA 등은 confidence x1.35, 엔터주 JYP·SM·HYBE·YG는 x1.3.\n"
        "DB 설계: sns_posts 테이블을 Alembic 마이그레이션으로 직접 작성·적용. ticker+platform 복합 인덱스와 posted_at 인덱스로 조회 최적화.\n"
        "표준 응답 DTO SnsSignalResult가 메인 에이전트의 confidence x source_tier_weight 합산에 자연스럽게 통합됩니다."
    ),
    38: CONT_NOTE,
    39: (
        "김세진 세 번째 슬라이드 — 자체 발견·수정한 버그와 Redis 캐시, 프론트엔드입니다.\n"
        "EUC-KR 디코딩 버그: 네이버 종목토론 한글이 깨져 DB에 저장되던 문제를, resp.content.decode 대신 BeautifulSoup에 from_encoding=euc-kr 직접 전달해 해결했습니다.\n"
        "DI 버그: 뉴스 라우터에서 AnalyzeNewsSignalUseCase 생성 시 keyword_resolver 파라미터가 누락된 것을 발견하고 직접 수정.\n"
        "Redis 캐시 + 자동 collect 트리거: analyze 요청 시 캐시 HIT이면 즉시 반환, DB 게시물 5건 미만이면 CollectSnsPostsUseCase 자동 호출 후 GPT 분석 진행.\n"
        "프론트엔드 SnsSignalCard: IDLE → LOADING → SUCCESS/ERROR 4개 상태. SUCCESS 시 시그널 배지·신뢰도 게이지·플랫폼별 결과·근거 게시물 top3·AI 한국어 요약으로 구성."
    ),
    40: CONT_NOTE,
    41: (
        "문우성 첫 번째 슬라이드 — 거시경제 현황판 데이터 수집 파이프라인입니다.\n"
        "매일 자동 수집 대상: 금리(기준금리·국채 10년물), 유가(WTI·브렌트), 환율(USD/KRW·DXY), 주요 경제 일정(FOMC·CPI·실업률 발표).\n"
        "소스별 수집 주기: FRED 12개 시리즈·유가·환율·VIX는 일 1회 / 경제 일정 크롤링은 주 1회.\n"
        "Redis 스냅샷 관리: _try_restore_macro_snapshot으로 4시간 이내 스냅샷이 Redis에 있으면 메모리 store에 복원합니다. 서버 재시작이나 hot-reload 시 LLM 호출 0회로 YouTube quota를 보호합니다."
    ),
    42: (
        "문우성 두 번째 슬라이드 — Risk-ON / Risk-OFF 자동 판단입니다.\n"
        "핵심 변수 4가지를 조합합니다: 금리 방향(상승 vs 하락), 달러 강도(DXY), VIX 수준(공포 지수), 유가 방향(에너지 비용).\n"
        "이 변수들을 LLM에 주입해 Risk-ON 또는 Risk-OFF를 결정하고, 근거 텍스트를 자동 생성합니다.\n"
        "예시: VIX 25 돌파 + 달러 강세 → Risk-OFF → 미 국채·금 선호, 위험자산 매도 압력 텍스트 자동 출력.\n"
        "/macro 페이지에서 금리·유가·환율 카드(전일 대비 등락 색상), 경제 일정 캘린더, Risk 배지+LLM 근거로 시각화합니다."
    ),
    43: (
        "문우성 세 번째 슬라이드 — 메인 에이전트 연결과 아키텍처 기여입니다.\n"
        "거시경제 에이전트가 MarketRiskJudgementResponse를 생성 → Redis 4h 캐시 저장 → 메인 에이전트 호출 시 거시 스냅샷으로 로드 → 종합 시그널 산출의 컨텍스트로 활용됩니다.\n"
        "거시 환경이 불안정할 때 전체 confidence를 하향 반영하는 구조입니다.\n"
        "아키텍처 기여: 전체 DDD·헥사고날 설계 참여, 레이어별 의존성 규칙 문서화, 도메인 간 인터페이스 표준화 방향 제시, 외부 API 공통 Client 패턴 정립."
    ),
    44: (
        "김민규 첫 번째 슬라이드 — History Agent 8단계 타임라인 파이프라인입니다.\n"
        "한 줄 요약: 종목 코드 하나로 8개 외부 소스에서 데이터를 수집해 시점이 명확한 사건만 골라 타임라인을 자동 생성합니다.\n"
        "Step 0~1: ticker 정규화 + 자산 유형 판별(EQUITY/INDEX/ETF) → Redis 캐시 HIT이면 즉시 반환(~50ms).\n"
        "Step 2: asyncio.gather로 yfinance / DART·SEC / Finnhub → GDELT → Naver 4경로 병렬 수집.\n"
        "Step 3~5: Jaccard >= 0.8 중복 제거 → DB enrichment 재사용 → 신규 이벤트 병렬 enrichment(Causality Agent 포함).\n"
        "Step 6~8: EventClassifier v2 + importance_score → DB upsert → Redis 3600s TTL 후 응답."
    ),
    45: CONT_NOTE,
    46: (
        "김민규 두 번째 슬라이드 — 5중 이상치 탐지기와 이벤트 분류 체계입니다.\n"
        "5중 탐지기: z-score(|수익률|>K×σ, 노랑) / cumulative 5d(±10% 진입, 오렌지) / cumulative 20d(±15%, 진홍) / drawdown(60봉 고점 대비 -10%→-3% 회복, 보라/에메랄드) / volatility cluster(5거래일 내 |r|>5% 2건 이상, 앰버).\n"
        "σ 계산 방식도 선택 가능합니다: stdev / stable / mad — settings.anomaly_robust_sigma_method 설정으로 운영 중 변경 가능.\n"
        "이벤트 분류 4카테고리: CORPORATE(yfinance 코드 매칭), ANNOUNCEMENT(SEC Item+LLM 재분류), MACRO(FRED+threshold), NEWS(source 필드 구분).\n"
        "MACRO TYPE_A(FOMC 발표 같은 원인 자체)와 TYPE_B(VIX 급등 같은 시장 반응, 추정 사유+신뢰도 포함)를 구분하는 것이 특히 중요합니다."
    ),
    47: (
        "김민규 세 번째 슬라이드 — Causality Agent, Dashboard UI, SSE 스트리밍입니다.\n"
        "Causality Agent: 이상치 봉 탐지 → 사용자가 마커 클릭 시 lazy 호출(LLM 비용 절약) → LangGraph StateGraph 4노드(수집→Tool Use→가설 생성→환각 방지) → 가설 3~6개+신뢰도 자동 생성 → DB 영구 저장. 두 번째 요청부터 ~50ms 응답.\n"
        "Dashboard UI: /dashboard에서 2열 4:1 그리드. NasdaqChart(이상치 마커 6종)+HistoryPanel(타임라인 카드) | StockSearch+AssetProfilePanel. Jotai 2.19 atomWithStorage로 localStorage 영속.\n"
        "SSE 스트리밍 API: GET /api/v1/history-agent/timeline/stream. progress/done/error 3종 이벤트로 4~8초 로딩 대신 실시간 진행 상황을 표시합니다."
    ),
    48: CONT_NOTE,
    49: (
        "이민영 첫 번째 슬라이드 — 관심종목 기반 개인화입니다.\n"
        "watchlist 관리 UI: 테마 그룹별 종목 버튼(반도체/자동차/바이오/엔터) 클릭으로 토글 등록·해제, 현재 선택 개수 실시간 표시.\n"
        "API: GET/POST/PUT/DELETE 4개 엔드포인트. 인증은 Redis 세션 기반(쿠키 user_token 또는 Authorization: Bearer 모두 지원).\n"
        "뉴스 개인화: 관심종목 없으면 전체 최신 뉴스 100건, 있으면 등록 종목명으로 DB 필터링(종목당 최대 10건, URL 중복 제거 후 최신순 정렬).\n"
        "유튜브 개인화: 관심종목 없으면 일반 주식 키워드 5개, 있으면 종목명 최대 5개를 YouTube 병렬 검색 키워드로 사용합니다."
    ),
    50: (
        "이민영 두 번째 슬라이드 — 국내 Smart Money(KRX 자동 수집)입니다.\n"
        "스마트머니란 외국인·기관·저명 투자자의 거래 흐름을 추적해 큰 돈이 어디로 이동하는지 파악하는 분석입니다.\n"
        "KRX 자동 수집: 매 영업일 16:30 KST(장 마감 30분 후) KOSPI+KOSDAQ 전 종목의 외국인/기관합계/개인 순매수거래량·대금 수집 후 DB 저장.\n"
        "집중도 점수: 외국인 점수(내 순매수 나누기 전체 최대값) + 기관 점수의 평균 x100. 외국인과 기관이 동시에 매수한 종목만 교집합으로 추출합니다."
    ),
    51: CONT_NOTE,
    52: (
        "이민영 세 번째 슬라이드 — 글로벌·국내 저명 투자자 포트폴리오입니다.\n"
        "글로벌 포트폴리오 — SEC 13F 19명: Warren Buffett·Michael Burry·George Soros 등 전설적 투자자부터 Ken Griffin·Renaissance Technologies 등 퀀트 운용사까지.\n"
        "수집 파이프라인: 분기별(2·5·8·11월 15일) SEC EDGAR 최신 13F-HR XML 파싱 → OpenFIGI API로 CUSIP → 티커 변환 → 전 분기와 비교 → 신규편입/비중확대/비중축소/청산 배지 자동 생성.\n"
        "국내 포트폴리오 — DART 대량보유보고 10명: 국민연금·미래에셋·삼성·박현주 등. 매월 1일 03:00 KST 자동 수집, 동시 10개 병렬 API 요청.\n"
        "변동 배지 4종(신규편입/비중확대/비중축소/청산)으로 한눈에 포트폴리오 변화를 확인할 수 있습니다."
    ),
    53: CONT_NOTE,
    54: PART_NOTE,
    55: (
        "GitHub·OKR·Slack 협업 프로세스입니다.\n"
        "GitHub 워크플로우: upstream(EDDI-RobotAcademy) → origin(개인 fork) → PR → Merge commit. main 직접 푸시 금지, squash 금지(원본 커밋 SHA 보존), Conventional Commits 준수.\n"
        "주요 PR 규모: 김세진 Sentiment 도메인 신설 PR #93이 +1,727줄·44파일, 자동 collect+Redis 캐시 PR #106이 +592줄.\n"
        "OKR 7개 Key Result 모두 달성 — 메인 에이전트 통합 / 한국+미국 / Source Tier / Sentiment / 거시 / 히스토리·인과관계 / Watchlist+Smart Money."
    ),
    56: (
        "Slack 협업 사례입니다.\n"
        "채널: #general #dev-backend #dev-frontend #design-architecture #bugs.\n"
        "실제 사례 — 영진님이 sentiment 빈 결과를 #bugs에 공유 → 세진님이 자동 collect + Redis 캐시 추가 커밋으로 즉시 해결. 실시간 소통이 품질을 높인 좋은 사례입니다."
    ),
    57: PART_NOTE,
    58: (
        "라이브 데모 시나리오입니다.\n"
        "① 삼성전자(005930): 5개 에이전트 병렬 분석, 통합 시그널+사업개요 카드, Source Tier 분포 확인.\n"
        "② AAPL: 영어 뉴스+yfinance 분기 EPS+SEC 8-K filing, 한국 종목과 동일 UI/UX로 표시.\n"
        "③ HYBE(352820): YouTube·X·Instagram 출처가 MEDIUM으로 자동 승격, 섹터 오버라이드 시그널 반영 확인.\n"
        "④ Smart Money 대시보드: 외국인+기관 집중 매수 카드, Warren Buffett 포트폴리오 변동 확인.\n"
        "⑤ 히스토리·인과관계: NASDAQ 차트+이상치 마커, 마커 클릭 → 원인 가설 자동 생성(LLM) 확인.\n"
        "⑥ Redis 분산 락 데모: 동일 종목 동시 5개 요청 → DART 1회만 호출됨을 로그 또는 LangSmith에서 확인."
    ),
    59: PART_NOTE,
    60: (
        "각 팀원이 이번 3·4차 프로젝트를 통해 배운 것과 개인 소감을 나누는 슬라이드입니다.\n"
        "김영진·김세진·문우성·김민규·이민영 5명의 개발 팀원, 이지은·정우진 2명의 외부 지원 팀원이 각자의 목소리로 이야기합니다.\n"
        "발표 당일 각 팀원이 직접 작성한 내용을 기반으로 발표합니다."
    ),
    61: PART_NOTE,
    62: (
        "프로젝트를 마무리하며 각 팀원이 하고 싶은 말을 전하는 슬라이드입니다.\n"
        "팀으로서 함께한 7인이 각자의 언어로 마무리 인사를 드립니다.\n"
        "발표 당일 각 팀원이 직접 작성한 내용을 기반으로 발표합니다."
    ),
    63: (
        "발표를 마치겠습니다. 감사합니다.\n"
        "저희 Antelligen의 핵심 메시지 — 투자 결과가 아니라, 검증 가능한 의사결정 과정과 안전한 실행 구조를 제공하는 서비스였습니다.\n"
        "GitHub 리포는 EDDI-RobotAcademy 조직에 공개되어 있습니다: antelligen-backend / antelligen-frontend.\n"
        "질문 있으시면 편하게 해주세요."
    ),
}


def update_notes(pptx_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    print(f"총 슬라이드: {total}")
    updated = 0
    for i, slide in enumerate(prs.slides, 1):
        note_text = NOTES.get(i)
        if note_text is None:
            continue
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        if tf is None:
            continue
        tf.clear()
        if note_text:
            tf.text = note_text
            updated += 1
    prs.save(str(pptx_path))
    print(f"업데이트된 슬라이드: {updated}개")
    print(f"저장 완료: {pptx_path}")


if __name__ == "__main__":
    update_notes(PPTX)
