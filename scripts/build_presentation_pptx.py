"""
PRESENTATION_SLIDES.md → PPTX 변환 스크립트.

사용법:
    .venv/Scripts/python.exe scripts/build_presentation_pptx.py

입력:  docs/angelligen/PRESENTATION_SLIDES.md
출력:  docs/angelligen/PRESENTATION_SLIDES.pptx
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
INPUT_MD = ROOT / "docs" / "angelligen" / "PRESENTATION_SLIDES.md"
OUTPUT_PPTX = ROOT / "docs" / "angelligen" / "PRESENTATION_SLIDES.pptx"

# 16:9 — 13.33in x 7.5in
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

CONTENT_LEFT = Inches(0.5)
CONTENT_RIGHT = Inches(12.83)  # SLIDE_WIDTH - 0.5
CONTENT_TOP = Inches(1.10)
CONTENT_BOTTOM = Inches(7.30)
CONTENT_WIDTH = Inches(12.33)

# Colors
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_ACCENT = RGBColor(0x29, 0x80, 0xB9)
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PART_BG = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_PART_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TABLE_HEADER_BG = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TABLE_ROW_ALT = RGBColor(0xF2, 0xF6, 0xFA)


# Speaker notes keyed by source-markdown slide number.
SPEAKER_NOTES: dict[int, str] = {
    1: (
        "안녕하세요, 저희 팀 Antelligen 발표를 시작하겠습니다.\n"
        "Antelligen은 AI 멀티 에이전트 기반 투자 인텔리전스 플랫폼입니다.\n"
        "핵심 메시지는 'AI가 대신 투자하는 시대가 아니라, AI의 판단을 사람이 검증하고 통제하는 구조'입니다."
    ),
    2: (
        "단일 LLM 기반 서비스는 복잡한 문제를 체계적으로 분해하지 못하고 한 모델에 과도하게 의존합니다.\n"
        "현실의 의사결정은 정보 수집·분석·판단·추천이 여러 전문가 협업으로 이루어집니다.\n"
        "그래서 역할별로 분리된 에이전트가 협업하는 구조가 필요했습니다."
    ),
    3: (
        "5가지 목표를 중심으로 설계했습니다 — 구조적 AI 시스템, 확장 가능한 아키텍처, 현실적 협업 흐름, 사용자 맞춤 분석, 자동화·효율성.\n"
        "결국 '여러 AI가 협업해서 문제를 해결하는 새로운 구조'를 만드는 게 우리 목표였습니다."
    ),
    4: (
        "글로벌 AI 에이전트 시장은 2025년 76억 달러에서 2033년 1,820억 달러로 연 50% 성장 전망입니다.\n"
        "국내 BFSI는 AI 도입이 가장 활발한 산업이고, 우리의 SOM은 중소기업과 개인 투자자입니다.\n"
        "클라우드 기반 웹서비스가 핵심 수익 모델이 됩니다."
    ),
    5: (
        "MAS를 도입하려는 고객은 세 가지 문제로 망설입니다 — 가시성·비용 통제·거버넌스.\n"
        "특히 금융에서는 '신뢰와 통제가 없는 AI는 실패한다'는 게 핵심 인사이트입니다.\n"
        "이 세 가지를 해결하는 게 Antelligen의 출발점이었습니다."
    ),
    6: (
        "데이터 신뢰도와 커버리지 문제 5가지를 해결했습니다.\n"
        "뉴스 8종목 제한을 풀고, 미국 주식까지 확장했고, 잠정실적까지 반영했습니다.\n"
        "출처별 가중치와 섹터 오버라이드로 신호 품질을 메타데이터로 표현했다는 점이 핵심입니다."
    ),
    7: (
        "운영·비용·UX 측면 6가지 문제도 정리했습니다 — quota·동시성·응답 지연·맥락 부족·부분 실패·Explainability.\n"
        "여기서 도출한 6대 해결 원칙이 저희 설계 철학입니다.\n"
        "특히 '부분 실패는 기능, 전체 실패는 버그'와 '거버넌스는 사전 설계' 두 원칙은 코드 전반에 일관되게 적용했습니다."
    ),
    8: (
        "기술 프레임워크 경쟁자는 CrewAI·AutoGen·LangGraph·LlamaIndex가 있고, 서비스형은 Lindy·AutoGPT가 있습니다.\n"
        "하지만 'AI 빌더'와 '트레이딩 도구'는 많아도, 검증 가능한 의사결정 과정을 제공하는 서비스는 거의 없습니다.\n"
        "이 빈 자리가 우리의 기회입니다."
    ),
    9: (
        "차별화는 신뢰성·데이터 통합·UX 세 축입니다.\n"
        "Tracing·Source Tier·Audit Trail로 거버넌스를 강화했고, 5개 멀티 에이전트로 한국·미국 데이터를 통합했습니다.\n"
        "결론은 '투자 결과가 아니라, 검증 가능한 의사결정 과정과 안전한 실행 구조를 제공하는 서비스가 시장 승자가 된다'입니다."
    ),
    10: (
        "1차 타겟은 '왜 떨어졌지?'를 검색하는 개인 투자자, 2차는 신입 애널리스트와 콘텐츠 크리에이터입니다.\n"
        "블룸버그는 월 200만 원이 넘지만, 한국 공시 중심으로 합리적 가격에 제공합니다.\n"
        "초기엔 웹, 다음 PWA, 마지막에 모바일 앱 순으로 진입합니다."
    ),
    11: (
        "KPI는 사용량·품질·비용 효율·사업성 네 축으로 봅니다.\n"
        "MAU 30~50명에서 1만 명까지, 응답 속도 8초에서 3초까지, 분석 비용 200원에서 50원까지가 목표입니다.\n"
        "1차 마일스톤은 유료 전환율 3%, 월 매출 100만 원입니다."
    ),
    12: (
        "서비스 확장은 다섯 방향입니다 — 시장·자산 확장, 기능 확장, 분석 모듈 확장, 수익 모델, 데이터 네트워크 효과.\n"
        "특히 데이터 네트워크 효과 — 사용자 행동 → 검증 사례 → 규칙 개선 → 정확도 향상 플라이휠입니다.\n"
        "1년만 운영하면 후발주자가 따라잡기 어려운 구조가 됩니다."
    ),
    13: (
        "백엔드는 Python 3.13 + FastAPI + PostgreSQL + Redis 조합입니다.\n"
        "AI는 OpenAI·Anthropic을 LangChain·LangGraph·LangSmith로 묶었고, 시장 데이터는 yfinance·pykrx 등을 씁니다.\n"
        "프론트엔드는 Next.js 16과 React 19 기반입니다."
    ),
    14: (
        "Clean Architecture 기반에 의존성은 항상 외부에서 내부로 향합니다.\n"
        "도메인 레이어는 순수 Python만 허용하고 FastAPI·SQLAlchemy·Redis 같은 외부 의존성은 일체 import 금지입니다.\n"
        "이 규칙을 CLAUDE.md로 강제하고 있습니다."
    ),
    15: (
        "총 24개 도메인으로 구성됐습니다 — 사용자·시장·콘텐츠·통합 에이전트·커뮤니티 영역으로 나눠집니다.\n"
        "외부 연동은 11종 — OpenAI·DART·SEC·yfinance·FRED·SerpAPI·YouTube 등입니다.\n"
        "각 도메인은 독립적으로 개발·교체 가능한 구조입니다."
    ),
    16: (
        "요청은 라우터에서 시작해서 Use Case와 Port를 거쳐 Repository나 External Client로 흘러갑니다.\n"
        "백그라운드 작업은 7종 — 공시 수집·뉴스 수집·KRX 순매수·NASDAQ·거시 리스크·SEC 13F·DART 등이 자동 실행됩니다.\n"
        "Redis 4시간 캐시로 매크로 스냅샷을 영속화해서 hot-reload 비용을 줄였습니다."
    ),
    17: (
        "메인 에이전트는 5개 서브 에이전트를 병렬로 호출합니다 — News·Disclosure·Finance·Sentiment·Macro.\n"
        "asyncio.gather로 부분 실패를 허용하고, Postgres 1시간 + Redis 7일 이중 캐시를 씁니다.\n"
        "최종 출력은 Bullish/Neutral/Bearish + Confidence입니다."
    ),
    18: (
        "출처 신뢰도는 4단계 — DART/SEC가 HIGH, 글로벌 IB가 MEDIUM, 국내 IB가 MEDIUM_LOW, SNS가 LOW입니다.\n"
        "엔터테인먼트 섹터는 SNS가 본질이라서 LOW를 MEDIUM으로 자동 승격하는 오버라이드를 넣었습니다.\n"
        "weighted_score = score × confidence × tier_multiplier로 계산합니다."
    ),
    19: (
        "캐시 TTL은 도메인 진실성을 코드로 표현한 겁니다 — 매크로 4시간, 회사 1일, 사업개요 7일.\n"
        "동시 요청은 Redis SET NX EX로 atomic 분산 락을 걸어 처리합니다.\n"
        "asyncio.gather(return_exceptions=True)로 1개 죽어도 나머지로 시그널을 산출합니다."
    ),
    20: (
        "팀은 7명입니다 — 김영진·김세진·문우성·김민규·이민영이 개발을 맡았고, 이지은·정우진은 외부 시점 피드백을 주셨습니다.\n"
        "각자 담당 영역과 핵심 작업, 테크 역량은 표에 정리되어 있습니다.\n"
        "GitHub Fork & PR 워크플로우와 CLAUDE.md로 도메인 규약을 공유했습니다."
    ),
    21: (
        "Github은 Fork & PR 워크플로우로, Merge commit으로 원본 SHA를 보존했습니다.\n"
        "OKR은 7개 KR을 모두 달성했습니다 — 통합 분석·미국 시장·Source Tier·Sentiment·거시·히스토리·Watchlist.\n"
        "Slack에서 영진님 sentiment 빈 결과 공유 → 김세진이 자동 collect 추가 커밋으로 즉시 해결한 사례가 있습니다."
    ),
    22: (
        "기능 4가지를 먼저 보여드리겠습니다 — 종합 종목 분석, 히스토리·인과관계 대시보드, Smart Money 자금 흐름, 글로벌 포트폴리오.\n"
        "Smart Money는 외국인·기관·개인 순매수와 13F 글로벌 19명, DART 국내 10명을 모두 추적합니다."
    ),
    23: (
        "나머지 5가지 기능 — Sentiment 감정 분석, 거시 경제 현황판, 관심 종목 개인화, 회사 프로필, 카카오 OAuth입니다.\n"
        "관심 종목을 등록하면 뉴스와 유튜브 페이지가 자동으로 개인화됩니다.\n"
        "회사 프로필은 ETF·INDEX도 동일한 카드 형식으로 제공됩니다."
    ),
    24: (
        "이제 라이브 데모로 넘어갑니다.\n"
        "삼성전자(KR), AAPL(US), HYBE(엔터 섹터 오버라이드)를 차례로 보여드리고, Smart Money와 히스토리 대시보드, 마지막으로 동시 요청에서의 분산 락까지 시연합니다.\n"
        "각 시연이 끝나면 다음으로 부드럽게 넘어가도록 진행하겠습니다."
    ),
    25: (
        "이번 프로젝트에서 각자 학습한 것과 소감을 공유합니다.\n"
        "각 팀원이 한마디씩 직접 말씀드립니다."
    ),
    26: (
        "마지막으로 각자 마무리 인사를 드립니다.\n"
        "긴 시간 들어주셔서 감사합니다."
    ),
    27: (
        "발표는 여기까지입니다. 질의응답 받겠습니다.\n"
        "GitHub 리포는 EDDI-RobotAcademy 조직에 backend와 frontend 모두 공개되어 있습니다."
    ),
}

PART_NOTE = (
    "잠시 다음 파트로 넘어가겠습니다."
)


@dataclass
class Block:
    kind: str  # 'h3' | 'p' | 'list' | 'code' | 'table' | 'quote'
    text: str = ""
    items: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


@dataclass
class Slide:
    number: int
    title: str
    part: str | None
    blocks: list[Block] = field(default_factory=list)


@dataclass
class PartCover:
    title: str


# ────────────────────────────────────────────────────────────────────────
# Markdown parsing
# ────────────────────────────────────────────────────────────────────────


def parse_md(text: str) -> tuple[list[object], dict]:
    lines = text.splitlines()
    items: list[object] = []
    current_part: str | None = None
    current_slide: Slide | None = None
    pending_part: PartCover | None = None
    i = 0
    metadata = {"title": "", "subtitle": ""}

    for line in lines:
        if line.startswith("# ") and not line.startswith("# Part"):
            metadata["title"] = line[2:].strip()
            break

    def flush_pending_part():
        nonlocal pending_part
        if pending_part is not None:
            items.append(pending_part)
            pending_part = None

    while i < len(lines):
        line = lines[i].rstrip()

        if line.startswith("# Part "):
            if current_slide is not None:
                items.append(current_slide)
                current_slide = None
            flush_pending_part()
            current_part = line[2:].strip()
            pending_part = PartCover(title=current_part)
            i += 1
            continue

        m = re.match(r"^## Slide\s+(\d+)\s*[—-]\s*(.+)$", line)
        if m:
            if current_slide is not None:
                items.append(current_slide)
            flush_pending_part()
            current_slide = Slide(
                number=int(m.group(1)),
                title=m.group(2).strip(),
                part=current_part,
            )
            i += 1
            continue

        if current_slide is None:
            i += 1
            continue

        if line.startswith("### "):
            current_slide.blocks.append(Block(kind="h3", text=line[4:].strip()))
            i += 1
            continue

        if line.startswith("```"):
            code_lines: list[str] = []
            i += 1
            while i < len(lines) and not lines[i].startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1
            # Strip leading/trailing blank lines but keep internal whitespace.
            while code_lines and code_lines[0].strip() == "":
                code_lines.pop(0)
            while code_lines and code_lines[-1].strip() == "":
                code_lines.pop()
            current_slide.blocks.append(Block(kind="code", text="\n".join(code_lines)))
            continue

        if line.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s:|-]+\|\s*$", lines[i + 1]):
            header = [c.strip() for c in line.strip("|").split("|")]
            i += 2
            rows: list[list[str]] = [header]
            while i < len(lines) and lines[i].startswith("|"):
                cells = [c.strip() for c in lines[i].strip("|").split("|")]
                rows.append(cells)
                i += 1
            current_slide.blocks.append(Block(kind="table", rows=rows))
            continue

        if line.startswith(">"):
            # Merge consecutive `>` lines into one multi-line quote block.
            quote_lines: list[str] = [line.lstrip("> ").rstrip()]
            i += 1
            while i < len(lines) and lines[i].startswith(">"):
                quote_lines.append(lines[i].lstrip("> ").rstrip())
                i += 1
            current_slide.blocks.append(Block(kind="quote", text="\n".join(quote_lines).strip()))
            continue

        if re.match(r"^[-*]\s+", line) or re.match(r"^\d+\.\s+", line):
            list_items: list[str] = []
            while i < len(lines) and (
                re.match(r"^[-*]\s+", lines[i]) or re.match(r"^\d+\.\s+", lines[i])
            ):
                list_items.append(re.sub(r"^([-*]|\d+\.)\s+", "", lines[i]).strip())
                i += 1
            current_slide.blocks.append(Block(kind="list", items=list_items))
            continue

        if line.strip() == "" or line.strip() == "---":
            i += 1
            continue

        current_slide.blocks.append(Block(kind="p", text=line.strip()))
        i += 1

    if current_slide is not None:
        items.append(current_slide)

    return items, metadata


def strip_markdown_inline(text: str) -> str:
    """Strip markdown inline markers, but preserve identifiers like
    `_SECTOR_OVERRIDE_` and `integrated_analysis_orm` (underscores between
    word chars), and drop dangling `**` markers from cross-line bold spans.
    """
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # Italics: _foo_ only when not adjacent to word chars on either side
    text = re.sub(r"(?<![\w_])_([^_\s][^_]*?[^_\s]|[^_\s])_(?![\w_])", r"\1", text)
    # Drop any leftover lone ** markers (cross-line bold)
    text = text.replace("**", "")
    return text


# ────────────────────────────────────────────────────────────────────────
# Height estimation
# ────────────────────────────────────────────────────────────────────────


def korean_weighted_len(text: str) -> int:
    """Cell-count for monospace layout planning.
    - Hangul (U+AC00–D7AF) and CJK ideographs render as 2 cells.
    - Box-drawing (U+2500–257F), block elements (U+2580–259F), and arrow
      glyphs render as 1 cell in Latin monospace fonts.
    - All other ASCII / Latin / punctuation render as 1 cell.
    """
    total = 0
    for c in text:
        cp = ord(c)
        # Hangul syllables / Jamo
        if 0xAC00 <= cp <= 0xD7AF or 0x1100 <= cp <= 0x11FF or 0x3130 <= cp <= 0x318F:
            total += 2
        # CJK unified ideographs
        elif 0x4E00 <= cp <= 0x9FFF:
            total += 2
        else:
            total += 1
    return total


def estimate_wrapped_lines(text: str, width_in: float, font_pt: int) -> int:
    """Estimate number of wrapped lines given width and font size."""
    # Rough widths: ASCII char ≈ 0.55 * font_pt / 72 inches
    ascii_w = 0.55 * font_pt / 72.0
    chars_per_line = max(10, int(width_in / ascii_w))
    total = 0
    for line in text.split("\n"):
        weighted = korean_weighted_len(line)
        if weighted == 0:
            total += 1
        else:
            total += max(1, (weighted + chars_per_line - 1) // chars_per_line)
    return total


def line_height(font_pt: int) -> float:
    """Pt → inches with 1.28 line spacing for readable, balanced layout."""
    return font_pt * 1.28 / 72.0


def estimate_block_height(block: Block, width_in: float = 12.0) -> float:
    """Returns estimated height in inches for a block, including padding.
    Slightly under-estimates so pagination doesn't push small content
    onto a near-empty page.
    """
    if block.kind == "h3":
        lines = estimate_wrapped_lines(block.text, width_in, 15)
        # Extra leading space so h3 visually separates from prior block.
        return line_height(15) * lines + 0.22
    if block.kind == "p":
        lines = estimate_wrapped_lines(block.text, width_in, 11)
        return line_height(11) * lines + 0.06
    if block.kind == "quote":
        lines = estimate_wrapped_lines(block.text, width_in - 0.4, 12)
        return line_height(12) * lines + 0.24
    if block.kind == "list":
        h = 0.0
        for item in block.items:
            lines = estimate_wrapped_lines(item, width_in - 0.3, 11)
            h += line_height(11) * lines + 0.04
        return h + 0.06
    if block.kind == "code":
        line_count = max(1, block.text.count("\n") + 1)
        # Empirical: PowerPoint renders Lucida/Consolas at ~1.30 line height.
        return line_count * (10 * 1.32 / 72.0) + 0.30
    if block.kind == "table":
        rows = block.rows
        if not rows:
            return 0.0
        n_cols = max(1, len(rows[0]))
        col_w = width_in / n_cols
        h = 0.0
        for r_idx, row in enumerate(rows):
            font_pt = 11 if r_idx == 0 else 10
            row_lines = 1
            for cell in row:
                cell_lines = estimate_wrapped_lines(cell, col_w - 0.15, font_pt)
                row_lines = max(row_lines, cell_lines)
            h += line_height(font_pt) * row_lines + 0.10
        return h + 0.06
    return 0.3


# ────────────────────────────────────────────────────────────────────────
# Rendering helpers
# ────────────────────────────────────────────────────────────────────────


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align=PP_ALIGN.LEFT,
    font_name: str = "Malgun Gothic",
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def set_notes(slide, text: str) -> None:
    """Insert speaker notes text onto a slide."""
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_part_cover(prs: Presentation, cover: PartCover):  # returns slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PART_BG
    bg.line.fill.background()

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.4), Inches(1.5), Inches(0.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    add_text_box(
        slide,
        Inches(1.0), Inches(2.5), Inches(11.3), Inches(0.9),
        cover.title,
        font_size=44, bold=True, color=COLOR_PART_FG,
    )
    return slide


def add_title_bar(slide, title: str, slide_number: int):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.95)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    add_text_box(
        slide, Inches(0.5), Inches(0.18), Inches(11.0), Inches(0.62),
        title,
        font_size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide, Inches(11.7), Inches(0.25), Inches(1.4), Inches(0.5),
        f"#{slide_number}",
        font_size=13, bold=False, color=RGBColor(0xCC, 0xDD, 0xEE),
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )


def render_block(slide, block: Block, top: float, height_in: float):
    """Render a single block at given top/height (inches as float)."""
    left = CONTENT_LEFT
    width = CONTENT_WIDTH
    top_emu = Inches(top)
    height_emu = Inches(height_in)

    if block.kind == "h3":
        tb = slide.shapes.add_textbox(left, top_emu, width, height_emu)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.10)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = strip_markdown_inline(block.text)
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = COLOR_ACCENT

    elif block.kind == "p":
        tb = slide.shapes.add_textbox(left, top_emu, width, height_emu)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = strip_markdown_inline(block.text)
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT

    elif block.kind == "quote":
        qbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top_emu, width, height_emu
        )
        qbox.fill.solid()
        qbox.fill.fore_color.rgb = RGBColor(0xEC, 0xF3, 0xFA)
        qbox.line.color.rgb = COLOR_ACCENT
        tf = qbox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.10)
        tf.margin_bottom = Inches(0.10)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for idx, qline in enumerate(block.text.split("\n")):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = strip_markdown_inline(qline)
            run.font.name = "Malgun Gothic"
            run.font.size = Pt(12)
            run.font.italic = True
            run.font.color.rgb = COLOR_PRIMARY

    elif block.kind == "list":
        tb = slide.shapes.add_textbox(left, top_emu, width, height_emu)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        for idx, item in enumerate(block.items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.line_spacing = 1.30
            run = p.add_run()
            run.text = "• " + strip_markdown_inline(item)
            run.font.name = "Malgun Gothic"
            run.font.size = Pt(11)
            run.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(4)

    elif block.kind == "code":
        cbox = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top_emu, width, height_emu
        )
        cbox.fill.solid()
        cbox.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
        cbox.line.fill.background()
        tf = cbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.10)
        tf.margin_bottom = Inches(0.10)
        # Detect box-drawing chars; if present, use a font with full coverage.
        has_box = any(0x2500 <= ord(c) <= 0x257F or 0x2580 <= ord(c) <= 0x259F or
                      ord(c) in (0x25BC, 0x25B6, 0x25C0, 0x25B2)  # ▼ ▶ ◀ ▲
                      for c in block.text)
        # Box-drawing diagrams: prefer Lucida Console — fixed-width, no
        # kerning, and full Unicode box-drawing coverage on all Windows.
        # Cascadia Mono has subtle kerning/ligatures that misalign ASCII art.
        code_font = "Lucida Console" if has_box else "Consolas"
        widest = max(
            (korean_weighted_len(line) for line in block.text.splitlines()),
            default=1,
        )
        usable_in = (width / 914400.0) - 0.40
        max_pt_for_width = max(7.0, usable_in / max(1, widest) / 0.0085)
        font_pt = max(7, min(10, int(max_pt_for_width)))
        for idx, line in enumerate(block.text.splitlines() or [""]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT  # paragraph[0] inherits shape default
            run = p.add_run()
            run.text = line
            run.font.name = code_font
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    elif block.kind == "table":
        rows = block.rows
        if not rows:
            return
        n_rows = len(rows)
        n_cols = len(rows[0])
        shape = slide.shapes.add_table(n_rows, n_cols, left, top_emu, width, height_emu)
        tbl = shape.table
        # Distribute total height across rows proportional to per-row content.
        # Header gets a fixed slot; remaining height divides among data rows
        # weighted by the max wrapped lines in that row.
        col_w_in = (width / 914400.0) / max(1, n_cols)
        per_row_weight: list[int] = []
        for r_idx, row in enumerate(rows):
            font_pt = 11 if r_idx == 0 else 10
            row_lines = 1
            for cell in row:
                cell_lines = estimate_wrapped_lines(cell, col_w_in - 0.20, font_pt)
                row_lines = max(row_lines, cell_lines)
            per_row_weight.append(row_lines)
        total_weight = sum(per_row_weight)
        total_h_emu = height_emu
        for r_idx in range(n_rows):
            row_h_emu = int(total_h_emu * per_row_weight[r_idx] / total_weight)
            tbl.rows[r_idx].height = row_h_emu
        for r_idx, row in enumerate(rows):
            for c_idx, cell_text in enumerate(row):
                if c_idx >= n_cols:
                    break
                cell = tbl.cell(r_idx, c_idx)
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.06)
                cell.margin_bottom = Inches(0.06)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text = ""
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = strip_markdown_inline(cell_text)
                run.font.name = "Malgun Gothic"
                if r_idx == 0:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TABLE_HEADER_FG
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_TABLE_HEADER_BG
                else:
                    run.font.size = Pt(10)
                    run.font.color.rgb = COLOR_TEXT
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        COLOR_TABLE_ROW_ALT if r_idx % 2 == 0 else COLOR_BG
                    )


# ────────────────────────────────────────────────────────────────────────
# Slide layout — distribute blocks; auto-paginate if overflow
# ────────────────────────────────────────────────────────────────────────


def split_into_pages(blocks: list[Block], available_h: float) -> list[list[tuple[Block, float]]]:
    """Greedy split with widow control and slack tolerance.

    - Allow content up to (available_h * SLACK) before forcing a page break.
      This lets slightly-overflowing content squeeze onto one slide rather
      than spawn a near-empty 2/2 page. The renderer will compress whitespace.
    """
    SLACK = 1.10  # tolerate 10% over available before paginating
    pages: list[list[tuple[Block, float]]] = []
    current: list[tuple[Block, float]] = []
    used = 0.0
    spacing = 0.08
    cap = available_h * SLACK

    def commit_page():
        nonlocal current, used
        if current and current[-1][0].kind == "h3" and len(current) > 1:
            widow = current.pop()
            used -= widow[1] + spacing
            pages.append(current)
            current = [widow]
            used = widow[1]
            return
        pages.append(current)
        current = []
        used = 0.0

    for block in blocks:
        h = estimate_block_height(block, width_in=12.0)
        if h > cap:
            if current:
                commit_page()
            pages.append([(block, available_h)])
            continue
        if used + h + (spacing if current else 0) > cap:
            commit_page()
            current = [(block, h)]
            used = h
        else:
            if current:
                used += spacing
            current.append((block, h))
            used += h
    if current:
        pages.append(current)
    return pages


def add_content_slide(prs: Presentation, slide_obj: Slide, page_blocks: list[tuple[Block, float]], page_idx: int, total_pages: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide_obj.title
    if total_pages > 1:
        title = f"{title} ({page_idx + 1}/{total_pages})"
    add_title_bar(slide, title, slide_obj.number)

    available_h = (float(CONTENT_BOTTOM) - float(CONTENT_TOP)) / 914400.0
    # Special case: single table on its own page → expand to fill height.
    if len(page_blocks) == 1 and page_blocks[0][0].kind == "table":
        block, _ = page_blocks[0]
        page_blocks = [(block, available_h)]

    base_spacing = 0.08
    used_blocks = sum(h for _, h in page_blocks)
    used_total = used_blocks + base_spacing * max(0, len(page_blocks) - 1)

    extra_gap = 0.0
    if used_total > available_h:
        # Slack overflow: compress block heights uniformly to fit available_h.
        scale = (available_h - base_spacing * max(0, len(page_blocks) - 1)) / used_blocks
        page_blocks = [(b, h * scale) for b, h in page_blocks]
    else:
        leftover = available_h - used_total
        n_gaps = max(1, len(page_blocks) - 1)
        # Sparse pages → distribute leftover space as gap between blocks.
        if leftover > 0.4:
            extra_gap = (leftover * 0.85) / n_gaps

    cursor = float(CONTENT_TOP) / 914400.0
    for idx, (block, h) in enumerate(page_blocks):
        if idx > 0:
            cursor += base_spacing + extra_gap
        render_block(slide, block, top=cursor, height_in=h)
        cursor += h
    return slide


def add_title_slide(prs: Presentation, slide_obj: Slide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    add_text_box(
        slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.1),
        "Antelligen",
        font_size=60, bold=True, color=COLOR_PART_FG,
    )
    add_text_box(
        slide, Inches(1.0), Inches(3.4), Inches(11.3), Inches(0.7),
        "AI Multi-Agent 기반 투자 인텔리전스 플랫폼",
        font_size=26, bold=False, color=RGBColor(0xCC, 0xDD, 0xEE),
    )
    add_text_box(
        slide, Inches(1.0), Inches(4.3), Inches(11.3), Inches(1.4),
        '"AI가 대신 투자하는 시대가 아니라,\nAI의 판단을 사람이 검증하고 통제하는 구조"',
        font_size=18, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF),
    )
    add_text_box(
        slide, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5),
        "발표일: 2026-05-02   |   팀: Antelligen (7인)",
        font_size=14, bold=False, color=RGBColor(0x88, 0xAA, 0xCC),
    )
    return slide


def add_thank_you_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    add_text_box(
        slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
        "Thank You",
        font_size=72, bold=True, color=COLOR_PART_FG, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.7),
        "Q & A",
        font_size=28, bold=False, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
        "GitHub: EDDI-RobotAcademy/antelligen-backend · antelligen-frontend",
        font_size=14, bold=False, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER,
    )
    return slide


# ────────────────────────────────────────────────────────────────────────
# Main
# ────────────────────────────────────────────────────────────────────────


def main():
    text = INPUT_MD.read_text(encoding="utf-8")
    items, _meta = parse_md(text)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    available_h = (float(CONTENT_BOTTOM) - float(CONTENT_TOP)) / 914400.0  # inches

    rendered = 0
    for item in items:
        if isinstance(item, PartCover):
            slide = add_part_cover(prs, item)
            if slide:
                set_notes(slide, PART_NOTE)
            rendered += 1
        elif isinstance(item, Slide):
            if item.number == 1:
                slide = add_title_slide(prs, item)
                set_notes(slide, SPEAKER_NOTES.get(1, ""))
                rendered += 1
            elif "Thank You" in item.title:
                continue
            else:
                pages = split_into_pages(item.blocks, available_h)
                if not pages:
                    pages = [[]]
                for p_idx, page in enumerate(pages):
                    slide = add_content_slide(prs, item, page, p_idx, len(pages))
                    if p_idx == 0:
                        set_notes(slide, SPEAKER_NOTES.get(item.number, ""))
                    rendered += 1

    # Always end with a Thank You slide.
    slide = add_thank_you_slide(prs)
    set_notes(slide, SPEAKER_NOTES.get(27, ""))
    rendered += 1

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"Generated: {OUTPUT_PPTX}")
    print(f"Total slides: {rendered}")


if __name__ == "__main__":
    main()
